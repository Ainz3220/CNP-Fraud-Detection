"""Plain-language rewrite. Match unique old phrases so every box updates."""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

PATH = r"c:\Codes\CNP-Fraud-Detection\CNP_Fraud_Detection_11slide_Visual_Final.pptx"

# slide index, unique old snippet, new text
PAIRS = [
    (0, "Card-Not-Present e-commerce fraud", "Online card fraud.\nTested with Mauritius examples."),
    (1, "CONTEXT", "THE PROBLEM"),
    (1, "CNP fraud is rare, so accuracy alone is not enough", "Fraud is rare. Accuracy is not enough."),
    (1, "rare-event, analyst-facing", "A model can look 99% correct and still miss most fraud."),
    (1, "fraud prevalence after CNP", "of payments are fraud"),
    (1, "fraud to legitimate", "one fraud for 81 genuine"),
    (1, "fraud cases in 294,854", "fraud cases in 294,854 payments"),
    (1, "The research gap", "Why this project"),
    (1, "PCA-anonymised", "Many public datasets hide what each field means."),
    (1, "minority class", "A high accuracy score can still miss fraud."),
    (1, "US-centric evidence", "US results may not work the same in Mauritius."),
    (1, "accurate enough to be useful", "Build a detector that is accurate, easy to review, and testable in Mauritius."),
    (2, "DATA AND FEATURES", "THE DATA"),
    (2, "Named behavioural features", "I use simple facts about each payment."),
    (2, "unambiguously online", "Only online shop types. No in-store card payments."),
    (2, "FINAL MODELLING SET", "AFTER FILTERING"),
    (2, "transactions\n3,576 fraud cases", "payments\n3,576 fraud cases\nabout 1 fraud in 81"),
    (2, "CNP filter", "Online only"),
    (2, "Chronological split", "Time split"),
    (2, "20% calibration", "60% train  |  20% tune  |  20% test"),
    (2, "7 named features:", "Seven facts: amount, hour, age, distance, amount z-score, shop type, gender"),
    (2, "7 FEATURE INPUTS", "SEVEN INPUTS"),
    (3, "METHOD", "THE MODELS"),
    (3, "Three models share one", "Three models. Same input."),
    (3, "separates model capacity", "I change the model. I keep the data and the test the same."),
    (3, "One transformed input", "Same features"),
    (3, "engineered feature vector", "All three models see the same numbers for each payment."),
    (3, "engineered\nfeatures", "same\nfeatures"),
    (3, "interpretable\nbaseline", "simple\nbaseline"),
    (3, "non-linear\nensemble", "many\nrules"),
    (3, "boosted\ndetector", "strongest\nmodel"),
    (3, "NON-LINEAR", "RULES"),
    (3, "No SMOTE", "No fake extra fraud rows."),
    (3, "Native class weighting", "Fraud cases get extra weight."),
    (3, "PR-AUC and F1 guide", "I compare F1 and PR-AUC, not accuracy."),
    (3, "Runtime verdicts", "Each score becomes: APPROVED, REVIEW, or BLOCKED"),
    (4, "strongest standalone", "XGBoost is best on its own."),
    (4, "58,971-row", "Tested on 58,971 later payments. Random Forest is close behind."),
    (4, "sensitivity-oriented", "Random Forest: 91.29% F1.\nLogistic Regression finds more odd cases, with more false alarms."),
    (4, "Accuracy is omitted", "All three look 99% accurate. That number is misleading because fraud is rare."),
    (4, "BEST STANDALONE", "BEST ALONE"),
    (5, "EXPLAINABILITY AND ENSEMBLE", "EXPLAIN AND COMBINE"),
    (5, "A single strong signal", "If one model sees high risk, we block."),
    (5, "SHAP explains the signal", "SHAP shows why. The final decision uses the strongest warning."),
    (5, "SHAP turns model output", "SHAP turns a score into reasons."),
    (5, "Risk factors push", "Risk reasons point to fraud."),
    (5, "Legitimacy factors push", "Safe reasons point to a genuine payment."),
    (5, "text explainer converts", "The app writes this in plain English."),
    (5, "Most consistently important", "Often important: unusual amount, distance, and time of day"),
    (5, "RF 8%", "RF 9%\nAPPROVED"),
    (5, "MOST SEVERE VERDICT", "STRICTEST RESULT\nFRAUD BLOCKED"),
    (5, "544/545", "544 of 545 fraud cases were blocked. Some genuine payments were also blocked. That is the cost of being careful."),
    (6, "PROTOTYPE AND CONTEXT", "THE APP"),
    (6, "analyst-facing workflow", "A working app, not only a notebook."),
    (6, "real application connects", "Predict. Show reasons. Keep a history."),
    (6, "MAURITIUS BENCHMARKS", "MAURITIUS EXAMPLES"),
    (6, "The useful surprise", "What we learned"),
    (6, "E2 changed from the intended", "E2 was meant to need review. After the online-only filter, the models approved it. SHAP showed why."),
    (6, "MUR input is supported", "The form takes MUR. The model uses USD. The rate is 49."),
    (6, "Two of three benchmarks", "These three examples are demos. They are not a full Mauritius study."),
    (7, "contribution is practical", "The system works.\nMauritius proof is still early."),
    (7, "transparent detection workflow", "Next step: real local payment data."),
    (7, "What this study adds", "What I built"),
    (7, "01  Named behavioural", "01  Simple named features\n02  Reasons in plain English\n03  A careful combined decision"),
    (7, "Next research", "What is next"),
    (7, "Real Mauritian transaction data", "Real Mauritius data\nMore time-based features\nBetter score cut-offs"),
    (8, "INFERENCE FLOW", "HOW IT WORKS"),
    (8, "How one transaction moves", "What happens to one payment"),
    (8, "same decision path links", "From the form to a decision a person can check."),
    (8, "Raw transaction", "Amount, time,\nplace, shop type"),
    (8, "Hour", "Add distance\nand amount z-score"),
    (8, "fraud probabilities", "LR, RF, XGB\neach give a %"),
    (8, "analyst message", "Reasons\nin plain words"),
    (8, "Approved", "Approve, review,\nor block"),
    (8, "DESIGN PRINCIPLE", "MAIN IDEA"),
    (8, "Traceability is the product", "Every decision can be checked. The score is not a black box."),
    (9, "Three research questions", "Three questions. Three short answers."),
    (9, "Evidence supports the prototype", "The app is strong. Claims about Mauritius stay careful."),
    (9, "Can supervised models detect", "Can models find rare online fraud?"),
    (9, "XGBoost achieved 92.63%", "Yes. XGBoost: 92.63% precision and 97.60% PR-AUC. Random Forest is close."),
    (9, "Can explanations and escalation", "Can we explain the decision and stay safe?"),
    (9, "SHAP exposes the drivers", "Yes. SHAP shows the reasons. The strictest model can still block."),
    (9, "Does the prototype generalise", "Does it already work for Mauritius?"),
    (9, "Two of three MUR benchmarks", "Too soon. Two of three demos matched. We still need real local data."),
    (9, "Accurate and explainable prototype", "Good detector. Easy to review. Local proof still needed."),
    (10, "Why three models, not one", "Why not one model?"),
    (10, "They fail on different cases. Escalation", "They miss different cases. We keep the strongest warning."),
    (10, "They fail differently", "They miss different cases"),
    (10, "One model would either miss", "One model can miss a new fraud pattern. Another can block too many genuine payments."),
    (10, "HOW COMBINED", "HOW WE COMBINE"),
    (10, "Most severe verdict wins", "We do not average"),
    (10, "Not an average. Any FRAUD", "If any model says BLOCK, the app blocks. Missing fraud costs more than an extra check."),
    (10, "WHAT YOU SAW", "LIVE DEMO"),
    (10, "Disagreement is the feature", "They disagreed. That helped."),
    (10, "MUR 500,000: LR 100%", "MUR 500,000: LR blocked at 100%. Trees approved. MUR 800 at 02:00 and 4000 miles: RF blocked at 43%."),
    (10, "DESIGN RULE", "SIMPLE RULE"),
    (10, "If any model says FRAUD BLOCKED", "If any model blocks, we block.\nIf one says review, we review.\nIf all approve, we approve."),
    (11, "Each model has a distinct job", "What each model is for"),
    (11, "Baseline sensitivity, main precision", "Three jobs. One careful final decision."),
    (11, "Linear weighted sum", "A simple straight-line model. A huge amount can push the score to 100%."),
    (11, "Sensitivity net", "Safety net"),
    (11, "Precision ~48%", "Precision about 48%  |  Recall about 90%"),
    (11, "flags out-of-range", "Good at extreme amounts. More false alarms."),
    (11, "200 shallow trees", "Many small rules. It looks at clues together, such as time and distance."),
    (11, "High-precision detector", "Everyday detector"),
    (11, "Precision ~90%", "Precision about 90%  |  Recall about 92%"),
    (11, "catches known fraud-like", "Good when several clues appear together. Weak on brand-new extremes."),
    (11, "Boosted trees", "Trees that learn from earlier mistakes. Best scores on the test set."),
    (11, "Best standalone scorer", "Best single model"),
    (11, "Precision ~92%", "Precision about 92%  |  PR-AUC about 98%"),
    (11, "highest ranking quality", "Best ranking of rare fraud. Same weak spot as Random Forest on brand-new extremes."),
    (12, "The explainer can contradict", "Trust the score. Check the red text."),
    (12, "Live check: prediction worked", "The % was right. Some warning words were extra rules, not the model."),
    (12, "PREDICTION - WORKING", "THE SCORE WORKED"),
    (12, "gauges match predict_proba", "The % comes from the model."),
    (12, "HIGH FRAUD text is not SHAP", "Red bullets are extra rules."),
    (12, "EXPLAINER - DOMAIN RULES", "THE TEXT CAN OVERSTATE"),
    (12, "Amount >= 5%", "Big amount always looks like risk.\nLate night always looks like risk.\nFar from home always looks like risk.\nEven if this model did not use it as fraud."),
    (12, "IF ASKED IN VIVA", "IF ASKED"),
    (12, "Trust the % and APPROVED / FRAUD", "Trust the % and APPROVED or BLOCKED.\nDo not treat HIGH FRAUD as the model's real reason."),
    (13, "Questions and discussion", "I am happy to take questions."),
]

NOTES = {
    0: "Say: this project finds online card fraud, and it shows why. Then walk through the data, the three models, the results, and the app.",
    1: "Say: fraud is only 1.21% of payments. That is one fraud in 81. So accuracy is a bad headline. I use precision, recall, F1 and PR-AUC. The aim is a detector people can review, and that we can test in Mauritius.",
    2: "Say: I kept only online shop types. 294,854 payments. Seven simple facts. Train on older payments. Test on later ones. Like real life.",
    3: "Say: same features go into three models. Logistic Regression is the simple baseline. Random Forest learns rules. XGBoost is the strongest. No fake fraud rows. Extra weight on real fraud. Final labels are approve, review, or block.",
    4: "Say: ignore 99% accuracy. XGBoost is best: about 93% precision and 98% PR-AUC. Random Forest is very close. Logistic Regression finds more odd cases but raises more false alarms.",
    5: "Say: SHAP shows which facts pushed the score. If any model blocks, we block. In the demo, LR was 100% and the trees were low. Combined still blocked. That catches more fraud. It also blocks some genuine payments.",
    6: "Say: this is a real app. E1 looks normal and is approved. E3 looks extreme and is blocked. E2 did not match the first plan after the online-only filter. SHAP showed why. These are demos, not a full Mauritius study.",
    7: "Say: I built a working, explainable detector. I have not proved Mauritius with real bank data yet. Next: real local data, more time features, and better cut-offs.",
    8: "Say: one payment goes through five steps. Capture. Add extra facts. Score with three models. Explain. Decide. A person can check every step.",
    9: "Say: question 1 yes, models can find rare fraud. Question 2 yes, we can show reasons and still block on the strongest warning. Question 3 too soon for Mauritius. We need real local data.",
    10: "Say: three models because they fail on different cases. We do not average. If any model blocks, we block. The 500,000 rupee case was caught by Logistic Regression. The late-night far-away small payment was caught by Random Forest.",
    11: "Say: Logistic Regression is the safety net for extreme amounts. Random Forest is the everyday model for mixed clues. XGBoost has the best test scores. I keep all three so one blind spot cannot hide fraud.",
    12: "Say: the percentages were correct. The red HIGH FRAUD lines are extra rules in the text, not always SHAP. Trust the score and the verdict. On the extreme case, XGBoost still approved, but Logistic Regression blocked, so the app blocked.",
    13: "Thank the audience. Invite questions. If needed, repeat: the system works, and Mauritius proof is still early.",
}


def copy_font(src_run, dest_run):
    dest_run.font.name = src_run.font.name
    dest_run.font.size = src_run.font.size
    dest_run.font.bold = src_run.font.bold
    dest_run.font.italic = src_run.font.italic
    try:
        dest_run.font.color.rgb = src_run.font.color.rgb
    except Exception:
        pass


def first_run(paragraph):
    for run in paragraph.runs:
        return run
    return None


def set_paragraph_text(paragraph, text, style_run):
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = text
        if style_run is not None:
            copy_font(style_run, run)
        return
    paragraph.runs[0].text = text
    if style_run is not None:
        copy_font(style_run, paragraph.runs[0])
    for extra in paragraph.runs[1:]:
        extra.text = ""


def apply_text(shape, text: str) -> None:
    tf = shape.text_frame
    lines = text.split("\n")
    style = None
    for p in tf.paragraphs:
        style = first_run(p)
        if style is not None:
            break
    for i, line in enumerate(lines):
        if i < len(tf.paragraphs):
            set_paragraph_text(tf.paragraphs[i], line, style)
        else:
            p = tf.add_paragraph()
            p.alignment = tf.paragraphs[0].alignment
            run = p.add_run()
            run.text = line
            if style is not None:
                copy_font(style, run)
    for j in range(len(lines), len(tf.paragraphs)):
        set_paragraph_text(tf.paragraphs[j], "", style)


def main():
    prs = Presentation(PATH)
    missing = []
    for slide_i, needle, new_text in PAIRS:
        if new_text == "SKIP":
            continue
        slide = prs.slides[slide_i]
        found = False
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            old = sh.text_frame.text.replace("\r", "")
            if needle in old:
                apply_text(sh, new_text)
                found = True
                break
        if not found:
            missing.append(f"slide {slide_i + 1}: {needle!r}")

    for slide_i, notes in NOTES.items():
        prs.slides[slide_i].notes_slide.notes_text_frame.text = notes

    prs.save(PATH)
    print("saved")
    if missing:
        print("MISSING", len(missing))
        for m in missing:
            print(" ", m)
    else:
        print("all pairs applied")


if __name__ == "__main__":
    main()

"""Simplify slide 6 to one idea: if any model blocks, we block."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import importlib.util

spec = importlib.util.spec_from_file_location(
    "simplify_pptx", r"c:\Codes\CNP-Fraud-Detection\scripts\simplify_pptx.py"
)
mod = importlib.util.module_from_spec(spec)
spec.loader.exec_module(mod)

PATH = r"c:\Codes\CNP-Fraud-Detection\CNP_Fraud_Detection_11slide_Visual_Final.pptx"

REMOVE = {
    "Oval 6",
    "Oval 8",
    "Oval 10",
    "Rectangle 7",
    "Rectangle 9",
    "Rectangle 11",
    "shap-features",
    "prediction-and-shap-interface-backing",
    "Picture 44",
    "ensemble-line",
    "visual-shap-flow-feature",
    "visual-shap-flow-arrow-1",
    "visual-shap-flow-impact",
    "visual-shap-flow-arrow-2",
    "visual-shap-flow-verdict",
}


def hide(shape):
    el = shape._element
    el.getparent().remove(el)


def place(shape, left, top, width=None, height=None):
    shape.left = int(left)
    shape.top = int(top)
    if width is not None:
        shape.width = int(width)
    if height is not None:
        shape.height = int(height)


def main():
    prs = Presentation(PATH)
    sl = prs.slides[5]
    by_name = {sh.name: sh for sh in sl.shapes}

    for name in list(REMOVE):
        sh = by_name.get(name)
        if sh is not None:
            hide(sh)

    by_name = {sh.name: sh for sh in sl.shapes}

    mod.apply_text(by_name["section-label"], "COMBINE")
    mod.apply_text(by_name["slide-title"], "If one model says block, we block.")
    mod.apply_text(by_name["slide-subtitle"], "We do not take a vote. We keep the strongest warning.")
    by_name["slide-subtitle"].height = Emu(400050)

    example = by_name["shap-title"]
    place(example, 533400, 1752600, 11125200, 457200)
    mod.apply_text(example, "Example: two models say approve. One says block. Final answer: block.")

    # Four big boxes in one row
    y = 2400300
    h = 1714500
    w = 2286000
    xs = [533400, 3238500, 5943600, 8648700]
    boxes = ["lr-output", "rf-output", "xgb-output", "ensemble-output"]
    texts = ["lr-output-text", "rf-output-text", "xgb-output-text", "ensemble-output-text"]
    labels = [
        "LR\n100%\nBLOCK",
        "RF\n9%\nAPPROVE",
        "XGB\n0%\nAPPROVE",
        "FINAL\nBLOCK",
    ]

    for name, x in zip(boxes, xs):
        place(by_name[name], x, y, w, h)
    for name, x, label in zip(texts, xs, labels):
        place(by_name[name], x + 95250, y + 228600, w - 190500, h - 381000)
        sh = by_name[name]
        tf = sh.text_frame
        tf.word_wrap = True
        try:
            tf._txBody.bodyPr.set("anchor", "ctr")
        except Exception:
            pass
        mod.apply_text(sh, label)
        for p in sh.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                if run.font.size is None or run.font.size < Pt(20):
                    run.font.size = Pt(22)

    arrows = ["ens-arrow-1", "ens-arrow-2", "ens-arrow-3"]
    arrow_x = [2966700, 5671800, 8376900]
    for name, x in zip(arrows, arrow_x):
        sh = by_name[name]
        place(sh, x, y + 685800, 266700, 342900)
        mod.apply_text(sh, "->")
        for p in sh.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.bold = True

    note = by_name["ensemble-note"]
    place(note, 533400, 4305300, 11125200, 571500)
    mod.apply_text(
        note,
        "This way we almost never miss fraud.\nWe also block some genuine payments. That is the cost of being careful.",
    )

    sl.notes_slide.notes_text_frame.text = (
        "Say only this. Point to the three boxes. Logistic Regression says 100 percent fraud. "
        "Random Forest and XGBoost say approve. We do not vote. We take the strictest answer. "
        "So we block. That catches almost all fraud. Some genuine payments get blocked too."
    )

    out = PATH
    try:
        prs.save(PATH)
    except PermissionError:
        out = r"c:\Codes\CNP-Fraud-Detection\CNP_Fraud_Detection_Presentation.pptx"
        prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    main()

"""Insert viva Q&A slides before the closing slide, matching deck chrome."""

from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(r"c:\Codes\CNP-Fraud-Detection")
SRC = ROOT / "CNP_Fraud_Detection_11slide_Visual_Final.pptx"
OUT = SRC
FALLBACK = ROOT / "CNP_Fraud_Detection_QandA.pptx"

NAVY = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x53, 0x61, 0x72)
LABEL = RGBColor(0x81, 0x90, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x5D, 0x58, 0xE9)
GREEN = RGBColor(0x27, 0x9A, 0x73)
ORANGE = RGBColor(0xE7, 0x86, 0x4F)
RED = RGBColor(0xC4, 0x45, 0x36)


def rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def duplicate_slide(prs: Presentation, index: int):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    sp_tree = dest.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}grpSp") or tag.endswith(
            "}cxnSp"
        ) or tag.endswith("}graphicFrame"):
            sp_tree.remove(child)

    src_cSld = source._element.cSld
    dst_cSld = dest._element.cSld
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        existing = dst_cSld.find(qn("p:bg"))
        if existing is not None:
            dst_cSld.remove(existing)
        dst_cSld.insert(0, deepcopy(src_bg))

    offset = 8000 + len(prs.slides) * 200
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        for el in new_el.iter():
            if el.tag.endswith("}cNvPr"):
                cur = el.get("id")
                if cur and cur.isdigit():
                    el.set("id", str(int(cur) + offset))
        sp_tree.append(new_el)

    return dest


def set_run(paragraph, text, size_pt, color, bold=False, italic=False):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def set_shape_text(shape, text, size_pt, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        set_run(p, line, size_pt, color, bold=bold, italic=italic)


def find_shape(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def strip_body(slide, keep_names: set[str]):
    sp_tree = slide.shapes._spTree
    for sh in list(slide.shapes):
        if sh.name in keep_names or sh.name.startswith("polish-progress-"):
            continue
        if sh.name.startswith("polish-title-"):
            continue
        sp_tree.remove(sh._element)


def add_rect(slide, left, top, width, height, fill, line=None, name="card"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Emu(9525)
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass
    return shape


def add_textbox(slide, left, top, width, height, text, size_pt, color, bold=False, align=PP_ALIGN.LEFT, name="text"):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", "t")
    except Exception:
        pass
    set_shape_text(box, text, size_pt, color, bold=bold, align=align)
    return box


def add_notes(slide, text: str):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def move_slide(prs: Presentation, old_index: int, new_index: int):
    sld_id_lst = prs.slides._sldIdLst
    items = list(sld_id_lst)
    el = items[old_index]
    sld_id_lst.remove(el)
    items = list(sld_id_lst)
    if new_index >= len(items):
        sld_id_lst.append(el)
    else:
        items[new_index].addprevious(el)


KEEP = {
    "rule",
    "section-label",
    "page-number",
    "slide-title",
    "slide-subtitle",
}


def prepare_qa_slide(prs, template_index, page, section, title, subtitle):
    slide = duplicate_slide(prs, template_index)
    strip_body(slide, KEEP)
    set_shape_text(find_shape(slide, "section-label"), section, 12, LABEL, bold=True)
    set_shape_text(find_shape(slide, "page-number"), page, 13, LABEL, align=PP_ALIGN.RIGHT)
    set_shape_text(find_shape(slide, "slide-title"), title, 26, NAVY, bold=True)
    sub = find_shape(slide, "slide-subtitle")
    sub.height = Emu(457200)
    set_shape_text(sub, subtitle, 14, MUTED)
    return slide


def build_why_three(slide):
    # Three evidence cards
    cards = [
        (
            "WHY THREE",
            "They fail differently",
            "One model would either miss new fraud patterns or block too many legitimate payments. The three classifiers cover different error modes.",
            "F7F9FB",
        ),
        (
            "HOW COMBINED",
            "Most severe verdict wins",
            "Not an average. Any FRAUD BLOCKED becomes the banner. Missing fraud costs more than an extra review.",
            "FFFFFF",
        ),
        (
            "WHAT YOU SAW",
            "Disagreement is the feature",
            "MUR 500,000: LR 100% blocked, trees approved. MUR 800 at 02:00 / 4000 miles: RF 43% blocked, LR 0.1% approved.",
            "F7F9FB",
        ),
    ]
    x0, y0, w, h, gap = 533400, 1785000, 3575050, 2857500, 200025
    for i, (kicker, head, body, fill) in enumerate(cards):
        x = x0 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill, "D6DDE5", name=f"why-card-{i}")
        add_textbox(slide, x + 171450, y0 + 152400, w - 342900, 190500, kicker, 11, PURPLE, bold=True, name=f"why-kicker-{i}")
        add_textbox(slide, x + 171450, y0 + 381000, w - 342900, 571500, head, 18, NAVY, bold=True, name=f"why-head-{i}")
        add_textbox(slide, x + 171450, y0 + 990600, w - 342900, 1619250, body, 14, MUTED, name=f"why-body-{i}")

    add_rect(slide, 533400, 4781550, 11125200, 1143000, "EAF5FB", "D6DDE5", name="why-callout")
    add_textbox(slide, 742950, 4933950, 1905000, 228600, "DESIGN RULE", 11, PURPLE, bold=True, name="why-callout-label")
    add_textbox(
        slide,
        2743200,
        4902200,
        8572500,
        857250,
        "If any model says FRAUD BLOCKED, the combined verdict is FRAUD BLOCKED.\nElse if any model says REVIEW REQUIRED, the combined verdict is REVIEW REQUIRED. Otherwise APPROVED.",
        14,
        NAVY,
        name="why-callout-text",
    )
    add_notes(
        slide,
        "Use this if asked why three models. The point is complementary failure modes, not three copies of the same score. Escalation is conservative on purpose: a missed fraud is more costly than a false block that a human can review.\n"
        "[Bullet points]\n"
        "- Not an average; most severe verdict wins\n"
        "- LR caught the MUR 500,000 extreme; RF caught 02:00 + 4000 miles on a small amount\n"
        "- Combined banner can block even when two models approve\n"
        "[Sources]\n"
        "- Live prototype tests; backend/models/predict.py majority_vote; saved_models/metrics.json",
    )


def build_each_model(slide):
    models = [
        (
            "BASELINE",
            "Logistic Regression",
            "Linear weighted sum. Extrapolates: a huge amount drives the score toward 100%.",
            "Sensitivity net",
            "Precision ~48%  |  Recall ~90%",
            "Advantage: flags out-of-range inputs the trees have never seen. Trade-off: more false positives.",
            "EAF5FB",
            "6DCBF4",
        ),
        (
            "MAIN",
            "Random Forest",
            "200 shallow trees. Learns combinations (hour x distance x category). Does not extrapolate.",
            "High-precision detector",
            "Precision ~90%  |  Recall ~92%",
            "Advantage: catches known fraud-like interactions. Trade-off: extreme amounts can stay in a mostly-legit leaf.",
            "F4F1FF",
            "5D58E9",
        ),
        (
            "ADVANCED",
            "XGBoost",
            "Boosted trees. Each tree tries to fix the previous errors. Strongest on patterns in the training data.",
            "Best standalone scorer",
            "Precision ~92%  |  PR-AUC ~98%",
            "Advantage: highest ranking quality under imbalance. Trade-off: same tree limit on brand-new extremes.",
            "F7F9FB",
            "279A73",
        ),
    ]
    x0, y0, w, h, gap = 533400, 1752600, 3575050, 4183380, 200025
    for i, (tier, name, how, role, metrics, adv, fill, accent) in enumerate(models):
        x = x0 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill, "D6DDE5", name=f"model-card-{i}")
        add_rect(slide, x, y0, Emu(57150), h, accent, name=f"model-rail-{i}")
        add_textbox(slide, x + 190500, y0 + 152400, w - 342900, 190500, tier, 11, rgb(accent) if accent != "6DCBF4" else PURPLE, bold=True, name=f"model-tier-{i}")
        add_textbox(slide, x + 190500, y0 + 381000, w - 342900, 495300, name, 18, NAVY, bold=True, name=f"model-name-{i}")
        add_textbox(slide, x + 190500, y0 + 914400, w - 342900, 857250, how, 13, MUTED, name=f"model-how-{i}")
        add_textbox(slide, x + 190500, y0 + 1828800, w - 342900, 266700, role, 13, NAVY, bold=True, name=f"model-role-{i}")
        add_textbox(slide, x + 190500, y0 + 2133600, w - 342900, 266700, metrics, 12, LABEL, bold=True, name=f"model-metrics-{i}")
        add_textbox(slide, x + 190500, y0 + 2476500, w - 342900, 1428750, adv, 13, MUTED, name=f"model-adv-{i}")
    add_notes(
        slide,
        "Walk the three roles. Logistic regression is the interpretable sensitivity net because a linear model extrapolates. Random Forest is the main high-precision detector for combinations. XGBoost is the strongest standalone model on held-out data. You would pick XGBoost alone for the best F1; you keep all three so a single blind spot cannot hide a high-risk case.\n"
        "[Bullet points]\n"
        "- LR: linear, extrapolates, high recall, lower precision\n"
        "- RF: interactions, high precision, weak on out-of-range amounts\n"
        "- XGB: best PR-AUC; same tree limit on extremes\n"
        "[Sources]\n"
        "- Dissertation model comparison; saved_models/metrics.json; live prototype screenshots",
    )


def build_explainer(slide):
    add_rect(slide, 533400, 1752600, 5334000, 2857500, "F7F9FB", "D6DDE5", name="pred-card")
    add_textbox(slide, 742950, 1879600, 4953000, 228600, "PREDICTION - WORKING", 12, GREEN, bold=True, name="pred-kicker")
    add_textbox(
        slide,
        742950,
        2171700,
        4953000,
        495300,
        "The gauges match predict_proba.",
        18,
        NAVY,
        bold=True,
        name="pred-title",
    )
    add_textbox(
        slide,
        742950,
        2762250,
        4953000,
        1524000,
        "MUR 500,000 at 02:00, 4000 miles\nLR 100%  FRAUD BLOCKED\nRF 9%  APPROVED\nXGB 0%  APPROVED\nCombined  FRAUD BLOCKED  (escalation)",
        14,
        MUTED,
        name="pred-body",
    )

    add_rect(slide, 6115050, 1752600, 5539725, 2857500, "FFFFFF", "D6DDE5", name="exp-card")
    add_textbox(slide, 6324600, 1879600, 5143500, 228600, "EXPLAINER - DOMAIN RULES", 12, RED, bold=True, name="exp-kicker")
    add_textbox(
        slide,
        6324600,
        2171700,
        5143500,
        495300,
        "HIGH FRAUD text is not SHAP.",
        18,
        NAVY,
        bold=True,
        name="exp-title",
    )
    add_textbox(
        slide,
        6324600,
        2762250,
        5143500,
        1619250,
        "Amount >= 5% above category mean -> always listed as risk\nHour outside 06:00-21:59 -> always listed as risk\nDistance >= 100 miles -> always listed as risk\nEven when that model's SHAP for the feature is negative.",
        14,
        MUTED,
        name="exp-body",
    )

    add_rect(slide, 533400, 4781550, 11125200, 1143000, "EAF5FB", "D6DDE5", name="exp-callout")
    add_textbox(slide, 742950, 4933950, 2286000, 228600, "IF ASKED IN VIVA", 11, PURPLE, bold=True, name="exp-callout-label")
    add_textbox(
        slide,
        3124200,
        4876800,
        8191500,
        914400,
        "Trust the % and APPROVED / FRAUD BLOCKED. Do not treat HIGH FRAUD INDICATOR as “why this model scored that.”\nOn the extreme case, XGBoost SHAP for amount, hour and distance was negative — it pushed toward legitimate — while the bullets still said HIGH FRAUD.",
        13,
        NAVY,
        name="exp-callout-text",
    )
    add_notes(
        slide,
        "If the demo shows 0% XGBoost next to three HIGH FRAUD bullets, the bug is the explainer overlay, not predict_proba. generate_explanation classifies amount, hour and distance with analyst rules even when SHAP is negative. Trees do not extrapolate, so MUR 500,000 can score near zero; LR hits 100%. Combined still blocks because LR blocked.\n"
        "[Bullet points]\n"
        "- Prediction path is correct; screenshots matched expected probabilities\n"
        "- Explainer ignores SHAP sign for amount, hour, distance\n"
        "- Empty age defaults to 40 — that is intended, not a data leak\n"
        "[Sources]\n"
        "- backend/utils/text_explainer.py; live SHAP dump for the MUR 500,000 case",
    )


def main():
    if not SRC.exists():
        raise SystemExit(f"Missing {SRC}")

    backup = ROOT / "CNP_Fraud_Detection_11slide_Visual_Final.before_qa.pptx"
    if not backup.exists():
        shutil.copy2(SRC, backup)

    prs = Presentation(str(SRC))
    template = 9  # existing ANSWERS slide

    s1 = prepare_qa_slide(
        prs,
        template,
        "11",
        "Q&A",
        "Why three models, not one",
        "They fail on different cases. Escalation keeps the most severe signal.",
    )
    build_why_three(s1)

    s2 = prepare_qa_slide(
        prs,
        template,
        "12",
        "Q&A",
        "Each model has a distinct job",
        "Baseline sensitivity, main precision, advanced ranking - then one conservative banner.",
    )
    build_each_model(s2)

    s3 = prepare_qa_slide(
        prs,
        template,
        "13",
        "Q&A",
        "The explainer can contradict the score",
        "Live check: prediction worked. HIGH FRAUD wording is a domain-rule overlay, not SHAP.",
    )
    build_explainer(s3)

    # New slides were appended after Thank You (old index 10). Move Thank You to the end.
    move_slide(prs, 10, 13)
    thanks = prs.slides[13]
    pn = find_shape(thanks, "thanks-page-number")
    if pn is not None:
        set_shape_text(pn, "14", 13, LABEL, align=PP_ALIGN.RIGHT)

    target = OUT
    try:
        prs.save(str(target))
        saved = target
    except OSError:
        prs.save(str(FALLBACK))
        saved = FALLBACK

    print(f"slides={len(prs.slides)} saved={saved}")
    for i, sl in enumerate(prs.slides, 1):
        title = ""
        for sh in sl.shapes:
            if sh.name in ("slide-title", "thanks-title", "cover-title", "closing-title"):
                title = sh.text_frame.text.replace("\n", " ")
                break
        print(f"  {i:02d}  {title}")


if __name__ == "__main__":
    main()

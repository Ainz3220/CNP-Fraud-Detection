"""Rebuild the 14-slide simple deck from the restored 11-slide original."""

from pptx import Presentation
import importlib.util

ROOT = r"c:\Codes\CNP-Fraud-Detection"


def load(name, path):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def main():
    add_qa = load("add_qa", ROOT + r"\scripts\add_qa_slides.py")
    simp = load("simp", ROOT + r"\scripts\simplify_pptx.py")
    s6 = load("s6", ROOT + r"\scripts\simplify_slide6.py")

    print("1 add Q&A slides")
    add_qa.main()
    print("2 simplify wording")
    simp.main()

    path = ROOT + r"\CNP_Fraud_Detection_11slide_Visual_Final.pptx"
    prs = Presentation(path)
    extra = [
        (2, "amount z-score, shop type", "Seven facts: amount, hour, age, distance, how unusual the amount is, shop type, gender"),
        (8, "and amount z-score", "Add distance\nand how unusual the amount is"),
        (9, "RESEARCH QUESTION", "QUESTION"),
        (9, "EVIDENCE / ANSWER", "ANSWER"),
        (12, "Combined  FRAUD BLOCKED", "MUR 500,000 at 02:00, 4000 miles\nLR 100%  BLOCKED\nRF 9%  APPROVED\nXGB 0%  APPROVED\nCombined  BLOCKED"),
        (6, "React dashboard", "Web dashboard"),
        (6, "-> FastAPI REST API", "-> API"),
        (6, "-> models + SHAP + SQLite history", "-> models, reasons, history"),
    ]
    for slide_i, needle, new_text in extra:
        for sh in prs.slides[slide_i].shapes:
            if sh.has_text_frame and needle in sh.text_frame.text.replace("\r", ""):
                simp.apply_text(sh, new_text)
                break
    prs.save(path)
    print("3 extra polish")

    print("4 simplify combine slide")
    s6.main()

    print("5 reorder")
    prs = Presentation(path)
    add_qa.move_slide(prs, 8, 6)
    add_qa.move_slide(prs, 8, 9)
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.name in ("page-number", "thanks-page-number") and sh.has_text_frame:
                simp.apply_text(sh, f"{i:02d}")
    prs.save(path)
    alt = ROOT + r"\CNP_Fraud_Detection_Presentation.pptx"
    prs.save(alt)
    print("saved", path)
    print("saved", alt)
    for i, sl in enumerate(prs.slides, 1):
        section = title = ""
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.replace("\n", " ").strip()
            if sh.name in ("section-label", "thanks-section-label"):
                section = t
            if sh.name in ("slide-title", "cover-title", "closing-title", "thanks-title"):
                title = t
        print(f"  {i:02d}  [{section or '-':20}]  {title}")


if __name__ == "__main__":
    main()
